"""Genera la presentación final en PowerPoint."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ──────────────────────────────────────────────────────────────────────────────
# Colores del tema
# ──────────────────────────────────────────────────────────────────────────────
AZUL_OSCURO  = RGBColor(0x0F, 0x29, 0x4D)   # fondo títulos
AZUL_MEDIO   = RGBColor(0x1A, 0x56, 0x9A)   # acentos
AZUL_CLARO   = RGBColor(0xD6, 0xE4, 0xF7)   # fondo suave
NARANJA      = RGBColor(0xE8, 0x6A, 0x1F)   # resaltado
BLANCO       = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_TEXTO   = RGBColor(0x2D, 0x2D, 0x2D)
GRIS_CLARO   = RGBColor(0xF4, 0xF6, 0xF8)

IMAGENES = Path("docs/imagenes")
ASSETS   = Path("presentacion_assets")

W = Inches(13.33)   # ancho slide widescreen
H = Inches(7.5)     # alto slide widescreen


# ──────────────────────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────────────────────

def _fondo(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _caja(slide, texto: str, x, y, w, h, *,
          size: int = 24, bold: bool = False, color: RGBColor = GRIS_TEXTO,
          align=PP_ALIGN.LEFT, fondo: RGBColor | None = None) -> None:
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if fondo:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = fondo


def _titulo_slide(slide, titulo: str, subtitulo: str = "") -> None:
    """Banda superior azul con título."""
    banda = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), W, Inches(1.4),
    )
    banda.fill.solid()
    banda.fill.fore_color.rgb = AZUL_OSCURO
    banda.line.fill.background()

    _caja(slide, titulo, Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.9),
          size=32, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)
    if subtitulo:
        _caja(slide, subtitulo, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.4),
              size=16, color=RGBColor(0xAD, 0xC8, 0xE8), align=PP_ALIGN.LEFT)


def _bala(slide, items: list[str], x, y, w, h, size: int = 20) -> None:
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"▸  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = GRIS_TEXTO


def _imagen(slide, ruta: Path, x, y, w=None, h=None) -> None:
    if not ruta.exists():
        return
    if w and h:
        slide.shapes.add_picture(str(ruta), x, y, w, h)
    elif w:
        slide.shapes.add_picture(str(ruta), x, y, width=w)
    elif h:
        slide.shapes.add_picture(str(ruta), x, y, height=h)
    else:
        slide.shapes.add_picture(str(ruta), x, y)


def _linea_decorativa(slide) -> None:
    """Línea naranja bajo la banda de título."""
    linea = slide.shapes.add_shape(1, Inches(0), Inches(1.4), W, Inches(0.06))
    linea.fill.solid()
    linea.fill.fore_color.rgb = NARANJA
    linea.line.fill.background()


def _numero(slide, n: int) -> None:
    _caja(slide, str(n), Inches(12.7), Inches(7.0), Inches(0.5), Inches(0.4),
          size=14, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.RIGHT)


# ──────────────────────────────────────────────────────────────────────────────
# Slides
# ──────────────────────────────────────────────────────────────────────────────

def slide_portada(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(slide, AZUL_OSCURO)

    # Franja decorativa izquierda
    franja = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.45), H)
    franja.fill.solid()
    franja.fill.fore_color.rgb = NARANJA
    franja.line.fill.background()

    _caja(slide, "Análisis y Diseño de Algoritmos",
          Inches(0.7), Inches(1.2), Inches(11), Inches(0.6),
          size=18, color=RGBColor(0xAD, 0xC8, 0xE8))

    _caja(slide, "Extensión de Bi-Particiones\na K-Particiones en IIT",
          Inches(0.7), Inches(1.9), Inches(11), Inches(2.0),
          size=44, bold=True, color=BLANCO)

    _caja(slide, "Búsqueda de la Partición de Mínima Información (MIP) con k ≥ 2 grupos",
          Inches(0.7), Inches(3.9), Inches(11), Inches(0.8),
          size=20, color=RGBColor(0xD0, 0xD8, 0xE8))

    _caja(slide, "ProyectoAnalisis2026  ·  2026",
          Inches(0.7), Inches(6.7), Inches(11), Inches(0.5),
          size=14, color=RGBColor(0x88, 0x99, 0xAA))


def slide_problema(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(slide, GRIS_CLARO)
    _titulo_slide(slide, "¿Cuál es el problema?", "Teoría de Información Integrada (IIT)")
    _linea_decorativa(slide)
    _numero(slide, 2)

    _bala(slide, [
        "IIT mide cuánta información integra un sistema (φ) que no puede recuperarse si se divide",
        "La Partición de Mínima Información (MIP) es el corte que menos destruye esa integración",
        "Encontrar la MIP es un problema de optimización combinatoria NP-difícil",
        "IIT clásica solo considera biparticiones (k=2): divide el sistema en exactamente 2 partes",
    ], Inches(0.5), Inches(1.6), Inches(7.5), Inches(4.5), size=21)

    _imagen(slide, IMAGENES / "01_biparticion.png",
            Inches(8.1), Inches(1.6), h=Inches(4.2))


def slide_extension(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(slide, GRIS_CLARO)
    _titulo_slide(slide, "Extensión a k-Particiones", "De 1 hiperplano de corte a k−1 hiperplanos")
    _linea_decorativa(slide)
    _numero(slide, 3)

    _caja(slide, "k=2  →  1 hiperplano de corte",
          Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
          size=18, bold=True, color=AZUL_MEDIO)
    _imagen(slide, ASSETS / "hipercubo_k2.png",
            Inches(0.5), Inches(2.1), w=Inches(5.8))

    _caja(slide, "k=3  →  2 hiperplanos de corte",
          Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5),
          size=18, bold=True, color=RGBColor(0x33, 0x88, 0x55))
    _imagen(slide, ASSETS / "hipercubo_k3.png",
            Inches(6.8), Inches(2.1), w=Inches(5.8))

    # Propiedad clave
    prop = slide.shapes.add_shape(1, Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.75))
    prop.fill.solid()
    prop.fill.fore_color.rgb = AZUL_OSCURO
    prop.line.fill.background()
    _caja(slide, "Propiedad garantizada:  φ(k) ≤ φ(k−1) ≤ … ≤ φ(2)  —  más grupos, igual o menor pérdida",
          Inches(0.6), Inches(6.5), Inches(12.2), Inches(0.65),
          size=17, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)


def slide_algoritmos(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(slide, GRIS_CLARO)
    _titulo_slide(slide, "Algoritmos Implementados", "Patrón Template Method + SA para sistemas grandes")
    _linea_decorativa(slide)
    _numero(slide, 4)

    # Tabla manual
    encabezados = ["Estrategia", "Enfoque", "k>2", "Complejidad"]
    filas = [
        ["Q-Nodos",      "Submodularidad + SA multi-arranque",       "✓", "O(n² log n) + SA"],
        ["Geometric",    "Hipercubo geométrico + DP de subconjuntos","✓", "O(n·2ⁿ) + SA"],
        ["Fuerza Bruta", "Enumeración completa de biparticiones",    "✗", "O(2ⁿ)"],
    ]

    col_x = [Inches(0.4), Inches(2.5), Inches(8.8), Inches(10.2)]
    col_w = [Inches(2.0), Inches(6.1), Inches(1.2), Inches(2.8)]
    row_y = [Inches(1.65), Inches(2.35), Inches(3.05), Inches(3.75)]
    row_h = Inches(0.6)

    for ci, enc in enumerate(encabezados):
        caja_h = slide.shapes.add_shape(1, col_x[ci], row_y[0], col_w[ci], row_h)
        caja_h.fill.solid()
        caja_h.fill.fore_color.rgb = AZUL_OSCURO
        caja_h.line.fill.background()
        _caja(slide, enc, col_x[ci] + Inches(0.05), row_y[0] + Inches(0.1),
              col_w[ci], row_h, size=16, bold=True, color=BLANCO)

    for ri, fila in enumerate(filas):
        bg = AZUL_CLARO if ri % 2 == 0 else BLANCO
        for ci, celda in enumerate(fila):
            caja_c = slide.shapes.add_shape(1, col_x[ci], row_y[ri + 1], col_w[ci], row_h)
            caja_c.fill.solid()
            caja_c.fill.fore_color.rgb = bg
            caja_c.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            color_txt = RGBColor(0x22, 0x88, 0x44) if celda == "✓" else (
                        RGBColor(0xCC, 0x33, 0x33) if celda == "✗" else GRIS_TEXTO)
            _caja(slide, celda, col_x[ci] + Inches(0.08), row_y[ri + 1] + Inches(0.1),
                  col_w[ci], row_h, size=15, color=color_txt)

    _bala(slide, [
        "Para n ≤ umbral → enumeración exacta (óptimo global garantizado)",
        "Para n grande  → BuscadorKRecocido: SA multi-cadena con criterio de Metropolis",
        "BuscadorKDP: inicialización con DP de subconjuntos O(3ⁿ·k) + refinamiento SA",
    ], Inches(0.5), Inches(4.6), Inches(12.5), Inches(2.5), size=18)


def slide_resultados(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(slide, GRIS_CLARO)
    _titulo_slide(slide, "Resultados Experimentales", "Q-Nodos vs Geometric — sistemas 10A a 25A")
    _linea_decorativa(slide)
    _numero(slide, 5)

    _imagen(slide, Path("resultados_25A") / "comparacion_Q_vs_G_k2.png",
            Inches(0.4), Inches(1.6), w=Inches(6.1))
    _imagen(slide, Path("resultados_25A") / "mip_confirmacion_k2.png",
            Inches(6.8), Inches(1.6), w=Inches(6.1))

    _bala(slide, [
        "Q-Nodos y Geometric coinciden en φ para k=2 en todos los sistemas pequeños",
        "Geometric es hasta 10× más rápido en sistemas de n>8 nodos",
        "k=3 mejora φ en sistemas con alta integración distribuida",
    ], Inches(0.5), Inches(6.0), Inches(12.5), Inches(1.3), size=17)


def slide_demo(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(slide, AZUL_OSCURO)

    franja = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.45), H)
    franja.fill.solid()
    franja.fill.fore_color.rgb = NARANJA
    franja.line.fill.background()

    _caja(slide, "Demo en Vivo",
          Inches(0.9), Inches(1.8), Inches(11), Inches(1.2),
          size=52, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)

    _bala_demo = [
        "Sistema A  ·  4 nodos  ·  estado inicial 1000",
        "Subsistema: alcance=1110, mecanismo=1110",
        "",
        "python demo_presentacion.py --paso 1",
        "  → Q-Nodos vs Geometric con k=2",
        "",
        "python demo_presentacion.py --paso 2",
        "  → Comparación k=2 vs k=3 en ambas estrategias",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.9), Inches(3.1), Inches(11), Inches(4.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(_bala_demo):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = item
        is_cmd = item.strip().startswith("python")
        is_arrow = item.strip().startswith("→")
        run.font.size = Pt(16 if (is_cmd or is_arrow) else 20)
        run.font.bold = is_cmd
        run.font.color.rgb = (NARANJA if is_cmd else
                              RGBColor(0xAD, 0xC8, 0xE8) if is_arrow else
                              RGBColor(0xDD, 0xEE, 0xFF))

    _numero(slide, 6)


def slide_conclusiones(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(slide, GRIS_CLARO)
    _titulo_slide(slide, "Conclusiones y Lecciones Aprendidas", "")
    _linea_decorativa(slide)
    _numero(slide, 7)

    _caja(slide, "Conclusiones",
          Inches(0.5), Inches(1.6), Inches(6.0), Inches(0.5),
          size=20, bold=True, color=AZUL_MEDIO)
    _bala(slide, [
        "φ(k) ≤ φ(2) siempre se cumple — más grupos no empeoran la solución",
        "Q-Nodos y Geometric coinciden en k=2 con la solución exacta (fuerza bruta)",
        "SA multi-cadena permite escalar a sistemas donde la enumeración es inviable",
        "DP de subconjuntos (BuscadorKDP) da mejor punto de inicio que SA puro",
    ], Inches(0.5), Inches(2.1), Inches(6.1), Inches(3.5), size=18)

    _caja(slide, "Lecciones Aprendidas",
          Inches(6.9), Inches(1.6), Inches(6.0), Inches(0.5),
          size=20, bold=True, color=AZUL_MEDIO)
    _bala(slide, [
        "La canonicalización de etiquetas es crítica para evitar duplicados en la búsqueda",
        "El umbral exacto/heurístico debe calibrarse por estrategia, no ser global",
        "Los tests de consistencia k=2 vs fuerza bruta son la mejor red de seguridad",
        "Trabajo futuro: métricas alternativas (Fisher-Rao, Wasserstein) y n > 25",
    ], Inches(6.9), Inches(2.1), Inches(6.1), Inches(3.5), size=18)

    # Línea final
    linea_f = slide.shapes.add_shape(1, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.05))
    linea_f.fill.solid()
    linea_f.fill.fore_color.rgb = NARANJA
    linea_f.line.fill.background()
    _caja(slide, "ProyectoAnalisis2026  ·  Análisis y Diseño de Algoritmos  ·  2026",
          Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.4),
          size=13, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────────────────────────────────────

def main() -> None:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_portada(prs)
    slide_problema(prs)
    slide_extension(prs)
    slide_algoritmos(prs)
    slide_resultados(prs)
    slide_demo(prs)
    slide_conclusiones(prs)

    salida = Path("presentacion_kparticiones.pptx")
    prs.save(salida)
    print(f"Presentación guardada en: {salida.resolve()}")
    print(f"Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
